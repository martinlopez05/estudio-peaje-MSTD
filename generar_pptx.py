from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# --- Colores ---
BG       = RGBColor(0x0f, 0x1b, 0x4c)
CARD     = RGBColor(0x1e, 0x2d, 0x6b)
YELLOW   = RGBColor(0xf5, 0xc8, 0x42)
WHITE    = RGBColor(0xff, 0xff, 0xff)
LBLUE    = RGBColor(0x8a, 0xb4, 0xf8)
BODY     = RGBColor(0xcc, 0xd6, 0xf6)
RED      = RGBColor(0xff, 0x6b, 0x6b)
GREEN    = RGBColor(0x69, 0xdb, 0x7c)
ACCENT   = RGBColor(0x3b, 0x5b, 0xdb)
DARKCARD = RGBColor(0x0a, 0x12, 0x35)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completamente en blanco

def add_slide():
    s = prs.slides.add_slide(blank)
    # fondo
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    # barra inferior degradada (simulada con rectángulo)
    bar = s.shapes.add_shape(1, 0, H - Pt(4), W, Pt(4))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return s

def txb(slide, text, x, y, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Segoe UI'
    return box

def rect(slide, x, y, w, h, fill=CARD, radius=False):
    shape = slide.shapes.add_shape(
        17 if radius else 1,   # 17=rounded rect, 1=rect
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = 0.05
    return shape

def card(slide, x, y, w, h, num, title, body, num_color=ACCENT):
    rect(slide, x, y, w, h, CARD)
    # número
    num_box = slide.shapes.add_shape(17, x + Inches(0.15), y + Inches(0.15),
                                     Inches(0.45), Inches(0.45))
    num_box.fill.solid(); num_box.fill.fore_color.rgb = num_color
    num_box.line.fill.background()
    tf = num_box.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = str(num); run.font.bold = True
    run.font.size = Pt(14); run.font.color.rgb = WHITE
    run.font.name = 'Segoe UI'
    # título
    txb(slide, title, x + Inches(0.7), y + Inches(0.18),
        w - Inches(0.8), Inches(0.3), size=11, bold=True, color=YELLOW)
    # cuerpo
    txb(slide, body, x + Inches(0.15), y + Inches(0.55),
        w - Inches(0.3), h - Inches(0.65), size=10, color=BODY)

def info_box(slide, x, y, w, h, title, body, title_color=YELLOW, body_size=11):
    rect(slide, x, y, w, h, CARD)
    txb(slide, title, x + Inches(0.2), y + Inches(0.15),
        w - Inches(0.3), Inches(0.25), size=10, bold=True, color=title_color)
    txb(slide, body, x + Inches(0.2), y + Inches(0.45),
        w - Inches(0.3), h - Inches(0.55), size=body_size, color=BODY)

def slide_title(slide, title, yellow_part, subtitle):
    txb(slide, title + ' ', Inches(0.6), Inches(0.3),
        Inches(12), Inches(0.8), size=36, bold=True, color=WHITE)
    txb(slide, yellow_part, Inches(0.6) + Pt(len(title) * 20), Inches(0.3),
        Inches(6), Inches(0.8), size=36, bold=True, color=YELLOW)
    txb(slide, subtitle, Inches(0.6), Inches(0.95),
        Inches(12), Inches(0.35), size=13, color=LBLUE)

def page_num(slide, n):
    txb(slide, str(n).zfill(2), W - Inches(0.8), H - Inches(0.4),
        Inches(0.6), Inches(0.3), size=12, bold=True, color=ACCENT,
        align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────
# SLIDE 1 — PORTADA
# ─────────────────────────────────────────────
s1 = add_slide()
txb(s1, 'MODELO, SIMULACIÓN Y TEORÍA DE LA DECISIÓN',
    Inches(0), Inches(1.1), W, Inches(0.4),
    size=13, color=LBLUE, align=PP_ALIGN.CENTER, bold=False)
txb(s1, 'Estudio de un', Inches(0), Inches(1.6), W, Inches(0.7),
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s1, 'Sistema de Peaje', Inches(0), Inches(2.2), W, Inches(0.9),
    size=54, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
txb(s1, 'Simulación de Colas con SimPy  ·  Análisis Económico',
    Inches(0), Inches(3.15), W, Inches(0.4),
    size=16, color=LBLUE, align=PP_ALIGN.CENTER)

# meta boxes
for i, (lbl, val) in enumerate([
    ('MATERIA', 'Modelo, Simulación y\nTeoría de la Decisión'),
    ('INTEGRANTES', 'Valentino Aimale\nMartin López'),
    ('PROFESORES', 'Esteban Gidekel\nGastón Lezcano · Guido Mehaudy'),
]):
    x = Inches(0.5) + i * Inches(4.1)
    rect(s1, x, Inches(4.1), Inches(3.9), Inches(1.5), CARD)
    txb(s1, lbl, x + Inches(0.2), Inches(4.2), Inches(3.5), Inches(0.3),
        size=9, bold=True, color=YELLOW)
    txb(s1, val, x + Inches(0.2), Inches(4.55), Inches(3.5), Inches(0.9),
        size=11, color=BODY)

# ─────────────────────────────────────────────
# SLIDE 2 — EL PROBLEMA
# ─────────────────────────────────────────────
s2 = add_slide()
txb(s2, 'El ', Inches(0.6), Inches(0.25), Inches(2), Inches(0.7),
    size=34, bold=True, color=WHITE)
txb(s2, 'Problema', Inches(1.05), Inches(0.25), Inches(4), Inches(0.7),
    size=34, bold=True, color=YELLOW)
txb(s2, 'a Resolver', Inches(3.3), Inches(0.25), Inches(3), Inches(0.7),
    size=34, bold=True, color=WHITE)
txb(s2, 'Autovía Olavarría – Azul  ·  Sistema de cobro de peaje',
    Inches(0.6), Inches(0.88), Inches(12), Inches(0.3), size=13, color=LBLUE)

info_box(s2, Inches(0.4), Inches(1.3), Inches(5.9), Inches(1.6),
         'EL SISTEMA',
         '· 2 estaciones: A (Ascendente) y D (Descendente)\n'
         '· 3 cabinas por estación, solo 1 abierta todo el día\n'
         '· 2 cabinas de 7–9h en A y 19–20h en D\n'
         '· 4 tipos de vehículos con distribuciones distintas')

info_box(s2, Inches(0.4), Inches(3.05), Inches(5.9), Inches(1.7),
         'TIEMPOS DE SERVICIO',
         '· Gran porte: Uniforme(45, 55) seg\n'
         '· Grandes / Motos: Exponencial (media = 30 seg)\n'
         '· Pequeños: Triangular (mín=15, moda=20, máx=35 seg)\n'
         '· Ambulancias / Policía: paso libre fuera de cola')

info_box(s2, Inches(6.6), Inches(1.3), Inches(6.3), Inches(3.45),
         'LA PREGUNTA CENTRAL',
         'La 3ra cabina cuesta $100 cada 10 minutos habilitarla.\n'
         'Solo puede abrirse en bloques mínimos de 10 minutos.\n\n'
         '¿Cuánto tiene que costar la multa por espera > 3 minutos\n'
         'para que al concesionario le convenga abrir la 3ra cabina?\n\n'
         'Demostrar con 95% de confianza.',
         body_size=12)
page_num(s2, 2)

# ─────────────────────────────────────────────
# SLIDE 3 — SUPUESTO (U, no ρ)
# ─────────────────────────────────────────────
s3 = add_slide()
txb(s3, 'El ', Inches(0.6), Inches(0.25), Inches(1.5), Inches(0.7),
    size=34, bold=True, color=WHITE)
txb(s3, 'Supuesto', Inches(0.98), Inches(0.25), Inches(4), Inches(0.7),
    size=34, bold=True, color=YELLOW)
txb(s3, 'Clave', Inches(3.6), Inches(0.25), Inches(3), Inches(0.7),
    size=34, bold=True, color=WHITE)
txb(s3, '¿Por qué los intervalos de llegada se duplicaron?  ·  U = Σ(λᵢ × E[Sᵢ]) / c',
    Inches(0.6), Inches(0.88), Inches(12), Inches(0.3), size=13, color=LBLUE)

info_box(s3, Inches(0.4), Inches(1.3), Inches(5.9), Inches(1.5),
         'FACTOR DE UTILIZACIÓN U',
         'U mide qué fracción del tiempo están ocupados los servidores.\n\n'
         'U < 1  →  sistema estable (la cola se drena)\n'
         'U ≥ 1  →  sistema explosivo (la cola crece sin límite)\n\n'
         'λᵢ = llegadas tipo i (veh/seg)  ·  E[Sᵢ] = servicio medio (seg)  ·  c = cabinas')

info_box(s3, Inches(0.4), Inches(2.95), Inches(5.9), Inches(2.2),
         'CÁLCULO CON TASAS ORIGINALES DEL ENUNCIADO',
         'No pico · c = 1:\n'
         '(1/60)×50 + (1/70)×30 + (1/40)×23.33 + (1/380)×30\n'
         '= 0.833 + 0.429 + 0.583 + 0.079  →  U = 1.924 / 1 = 1.92  ✗\n\n'
         'Pico · c = 2:\n'
         '(1/30)×50 + (1/40)×30 + (1/25)×23.33 + (1/380)×30\n'
         '= 1.667 + 0.750 + 0.933 + 0.079  →  U = 3.429 / 2 = 1.71  ✗')

# tabla de U
headers = ['Escenario', 'Cabinas', 'U', 'Estado']
rows = [
    ('Original · no pico', '1', '1.92', 'Explosivo ✗'),
    ('Original · pico',    '2', '1.71', 'Explosivo ✗'),
    ('Original · pico',    '3', '1.14', 'Explosivo ✗'),
    ('50/50 · no pico',    '1', '0.96', 'Estable ✓'),
    ('50/50 · pico',       '2', '0.86', 'Estable ✓'),
    ('50/50 · pico + 3ra', '3', '0.57', 'Estable ✓'),
]
col_w = [Inches(2.4), Inches(0.8), Inches(0.7), Inches(1.15)]
col_x = [Inches(6.6), Inches(9.05), Inches(9.9), Inches(10.65)]
row_h = Inches(0.38)
for ci, (cx, cw) in enumerate(zip(col_x, col_w)):
    hdr = rect(s3, cx, Inches(1.3), cw, row_h, ACCENT)
    txb(s3, headers[ci], cx + Inches(0.05), Inches(1.34),
        cw, row_h, size=10, bold=True, color=WHITE)
for ri, row in enumerate(rows):
    y = Inches(1.3) + row_h + ri * row_h
    bg = CARD if ri < 3 else RGBColor(0x1a, 0x35, 0x1a)
    for ci, (cx, cw) in enumerate(zip(col_x, col_w)):
        rect(s3, cx, y, cw, row_h, bg)
        col = RED if ri < 3 else GREEN
        c_use = col if ci >= 2 else (YELLOW if ri >= 3 else BODY)
        txb(s3, row[ci], cx + Inches(0.05), y + Inches(0.05),
            cw, row_h, size=10, color=c_use)

info_box(s3, Inches(6.6), Inches(4.0), Inches(6.3), Inches(1.15),
         'DISTRIBUCIÓN DEL TRÁFICO (SUPUESTO 50/50)',
         'La autovía es bidireccional. Con reparto 50/50 cada estación recibe\n'
         'la mitad del flujo → el intervalo entre llegadas se duplica.\n'
         'Sin este ajuste U > 1 en todos los casos: simulación inválida.',
         body_size=10)
page_num(s3, 3)

# ─────────────────────────────────────────────
# SLIDE 4 — FUNCIONES DEL CÓDIGO
# ─────────────────────────────────────────────
s4 = add_slide()
txb(s4, 'Funciones del ', Inches(0.6), Inches(0.25), Inches(5), Inches(0.7),
    size=34, bold=True, color=WHITE)
txb(s4, 'Código', Inches(4.35), Inches(0.25), Inches(4), Inches(0.7),
    size=34, bold=True, color=YELLOW)
txb(s4, 'Arquitectura de la simulación en SimPy',
    Inches(0.6), Inches(0.88), Inches(12), Inches(0.3), size=13, color=LBLUE)

funcs = [
    (1, 'PeajeConfig',          'Diccionario con intervalos de llegada (seg) por tipo de vehículo y estado (pico/no pico).'),
    (2, 'tiempo_servicio()',     'Genera segundos de atención según la distribución: Uniforme, Exponencial o Triangular.'),
    (3, 'Estacion',             'Clase que modela cada peaje: 3 cabinas como recursos SimPy y estado de cabinas activas.'),
    (4, 'gestor_horario_base()', 'Proceso que ajusta cuántas cabinas base están abiertas según la hora (1 normal, 2 en pico).'),
    (5, 'supervisor_dinamico()', 'Cada 10 min calcula la espera media. Si supera 3 min abre la 3ra cabina y registra $100.'),
    (6, 'generador_vehiculos()', 'Genera vehículos con llegadas exponenciales. 4 generadores en paralelo por estación.'),
    (7, 'vehiculo()',            'El vehículo elige la cabina con menor cola, espera, registra su espera y se atiende.'),
    (8, 'correr_simulacion()',   'Construye el entorno, lanza todos los procesos y ejecuta 24hs. Devuelve multados y costos.'),
    (9, 'main()',                'Corre 30 réplicas, calcula medias, IC al 95% con t de Student y la multa de equilibrio.'),
]

cw = Inches(4.1); ch = Inches(1.0)
for i, (num, title, body) in enumerate(funcs):
    col = i % 3; row = i // 3
    x = Inches(0.4) + col * (cw + Inches(0.1))
    y = Inches(1.3) + row * (ch + Inches(0.08))
    card(s4, x, y, cw, ch, num, title, body)

page_num(s4, 4)

# ─────────────────────────────────────────────
# SLIDE 5 — PASOS MONTE CARLO
# ─────────────────────────────────────────────
s5 = add_slide()
txb(s5, 'Los 5 Pasos de ', Inches(0.6), Inches(0.25), Inches(5.5), Inches(0.7),
    size=34, bold=True, color=WHITE)
txb(s5, 'Monte Carlo', Inches(5.05), Inches(0.25), Inches(5), Inches(0.7),
    size=34, bold=True, color=YELLOW)
txb(s5, 'Cómo se aplican en esta simulación',
    Inches(0.6), Inches(0.88), Inches(12), Inches(0.3), size=13, color=LBLUE)

pasos = [
    ('PASO 0', YELLOW,  BG,    'Construir el Modelo',
     'Clases Estacion, procesos gestor_horario_base, supervisor_dinamico, generador_vehiculos y vehiculo. '
     'Se definen entidades, recursos, lógica de colas y condiciones de apertura de cabinas.'),
    ('PASO 1', ACCENT,  WHITE, 'Generar Números Aleatorios',
     'El módulo random de Python genera U(0,1) internamente en cada llamada a '
     'random.expovariate(), random.uniform() y random.triangular().'),
    ('PASO 2', ACCENT,  WHITE, 'Transformar en Entradas Válidas',
     'Los U(0,1) se convierten en tiempos con las distribuciones reales: '
     'Uniforme(45,55) · Exponencial(λ=1/30) · Triangular(15,20,35) · llegadas Exponencial(λ=1/tasa).'),
    ('PASO 3', ACCENT,  WHITE, 'Ejecutar el Modelo',
     'env.run(until=86400) procesa 24hs de simulación. Cada vehículo llega, espera, '
     'se atiende y registra su espera. Se repite 30 veces por escenario.'),
    ('PASO 4', ACCENT,  WHITE, 'Contabilidad y Estadística',
     'Se cuentan vehículos con espera > 3 min. Con las 30 réplicas se calcula media, '
     'IC al 95% (t de Student, df=29) y la multa de equilibrio.'),
]

for i, (badge, badge_bg, badge_fg, title, body) in enumerate(pasos):
    y = Inches(1.3) + i * Inches(1.12)
    # badge
    b = rect(s5, Inches(0.4), y, Inches(1.1), Inches(0.85), badge_bg)
    txb(s5, badge, Inches(0.4), y + Inches(0.25), Inches(1.1), Inches(0.4),
        size=12, bold=True, color=badge_fg, align=PP_ALIGN.CENTER)
    # contenido
    rect(s5, Inches(1.6), y, Inches(11.3), Inches(0.85), CARD)
    txb(s5, title, Inches(1.75), y + Inches(0.08), Inches(4), Inches(0.3),
        size=11, bold=True, color=YELLOW)
    txb(s5, body, Inches(1.75), y + Inches(0.38), Inches(11.0), Inches(0.42),
        size=10, color=BODY)

page_num(s5, 5)

# ─────────────────────────────────────────────
# SLIDE 6 — SEMILLA
# ─────────────────────────────────────────────
s6 = add_slide()
txb(s6, 'Semilla para ', Inches(0.6), Inches(0.25), Inches(4.5), Inches(0.7),
    size=34, bold=True, color=WHITE)
txb(s6, 'Reproducibilidad', Inches(4.0), Inches(0.25), Inches(6), Inches(0.7),
    size=34, bold=True, color=YELLOW)
txb(s6, 'Control de la generación de números aleatorios  ·  random.seed()',
    Inches(0.6), Inches(0.88), Inches(12), Inches(0.3), size=13, color=LBLUE)

info_box(s6, Inches(0.4), Inches(1.3), Inches(5.9), Inches(1.35),
         '¿QUÉ ES LA SEMILLA?',
         'El módulo random genera una secuencia de U(0,1) a partir de un valor inicial '
         'llamado semilla. La misma semilla siempre produce la misma secuencia, sin '
         'importar qué distribución la consuma (exponencial, uniforme, triangular).')

rect(s6, Inches(0.4), Inches(2.75), Inches(5.9), Inches(2.35), DARKCARD)
txb(s6, 'IMPLEMENTACIÓN EN EL CÓDIGO',
    Inches(0.6), Inches(2.85), Inches(5.5), Inches(0.3),
    size=10, bold=True, color=YELLOW)
code_lines = [
    'for i in range(n_replicas):',
    '    random.seed(i)   # semilla para escenario base',
    '    m_base, _ = correr_simulacion(False)',
    '',
    '    random.seed(i)   # misma semilla para dinámico',
    '    m_din, c_din = correr_simulacion(True)',
]
for li, line in enumerate(code_lines):
    txb(s6, line, Inches(0.6), Inches(3.2) + li * Inches(0.27),
        Inches(5.5), Inches(0.27), size=9.5, color=LBLUE)

ventajas = [
    (1, 'Reproducibilidad',
     'Dos corridas distintas del programa producen exactamente los mismos resultados. '
     'Permite verificar y compartir resultados sin ambigüedad.'),
    (2, 'Comparación justa entre escenarios',
     'Al usar la misma semilla i para el escenario base y el dinámico en cada réplica, '
     'ambos reciben la misma secuencia de eventos. La diferencia se debe solo al escenario.'),
    (3, 'Reducción de varianza',
     'Comparar con los mismos números aleatorios reduce las diferencias entre réplicas, '
     'haciendo los intervalos de confianza más precisos.'),
]
for i, (num, title, body) in enumerate(ventajas):
    y = Inches(1.3) + i * Inches(1.4)
    card(s6, Inches(6.6), y, Inches(6.3), Inches(1.28), num, title, body)

page_num(s6, 6)

# ─────────────────────────────────────────────
# SLIDE 7 — RESULTADOS Y CONCLUSIÓN
# ─────────────────────────────────────────────
s7 = add_slide()
txb(s7, 'Resultados', Inches(0.6), Inches(0.25), Inches(4.5), Inches(0.7),
    size=34, bold=True, color=YELLOW)
txb(s7, ' y Conclusión', Inches(3.95), Inches(0.25), Inches(5), Inches(0.7),
    size=34, bold=True, color=WHITE)
txb(s7, '30 réplicas  ·  semilla = i por réplica  ·  IC 95% con t de Student (df = 29)',
    Inches(0.6), Inches(0.88), Inches(12), Inches(0.3), size=13, color=LBLUE)

# 3 cajas de resultados
res_data = [
    ('ESCENARIO BASE (sin 3ra cabina)',    '3.044', 'vehículos/día > 3 min\nIC 95%: (2.865 – 3.223)', WHITE),
    ('ESCENARIO DINÁMICO (con 3ra cabina)','1.073', 'vehículos/día > 3 min\nIC 95%: (1.034 – 1.112)', WHITE),
    ('COSTO EXTRA POR ABRIR 3RA CABINA',  '$4.693','promedio diario\nbloques de $100 / 10 min',      WHITE),
]
rw = Inches(4.1)
for i, (title, val, sub, vc) in enumerate(res_data):
    x = Inches(0.4) + i * (rw + Inches(0.1))
    rect(s7, x, Inches(1.3), rw, Inches(1.6), CARD)
    txb(s7, title, x + Inches(0.15), Inches(1.38), rw - Inches(0.2), Inches(0.3),
        size=9, bold=True, color=LBLUE)
    txb(s7, val,   x + Inches(0.15), Inches(1.7),  rw - Inches(0.2), Inches(0.55),
        size=28, bold=True, color=YELLOW)
    txb(s7, sub,   x + Inches(0.15), Inches(2.25), rw - Inches(0.2), Inches(0.55),
        size=10, color=BODY)

# fórmula
rect(s7, Inches(0.4), Inches(3.05), Inches(12.5), Inches(0.55), DARKCARD)
txb(s7, 'Multa de equilibrio  =  Costo cabina / Vehículos salvados  =  $4.693 / 1.971  =  $2,38 por vehículo',
    Inches(0.6), Inches(3.1), Inches(12.1), Inches(0.45),
    size=14, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# 2 cajas finales
rect(s7, Inches(0.4), Inches(3.75), Inches(6.2), Inches(2.8), RGBColor(0x1a, 0x35, 0x1a))
txb(s7, 'RESPUESTA AL PROBLEMA',
    Inches(0.6), Inches(3.83), Inches(5.8), Inches(0.3),
    size=10, bold=True, color=GREEN)
txb(s7, '$2,38', Inches(0.6), Inches(4.15), Inches(5.8), Inches(0.8),
    size=42, bold=True, color=GREEN)
txb(s7, 'Si la multa > $2,38  →  conviene abrir la 3ra cabina.\n'
        'Si la multa < $2,38  →  es más barato pagar las multas.\n'
        'Los IC no se superponen → diferencia significativa al 95%.',
    Inches(0.6), Inches(4.95), Inches(5.8), Inches(0.8),
    size=10.5, color=BODY)

rect(s7, Inches(6.8), Inches(3.75), Inches(6.1), Inches(2.8), CARD)
txb(s7, 'INTERPRETACIÓN',
    Inches(7.0), Inches(3.83), Inches(5.7), Inches(0.3),
    size=10, bold=True, color=LBLUE)
txb(s7, 'La 3ra cabina salva ~1.971 vehículos/día de exceder los 3 minutos.\n\n'
        'El costo de abrirla es ~$4.693/día.\n\n'
        'Con una multa de $2,38/vehículo, ambas opciones cuestan lo mismo.\n\n'
        'Por encima de ese valor, la cabina extra se paga sola.',
    Inches(7.0), Inches(4.2), Inches(5.7), Inches(2.1),
    size=11, color=BODY)

page_num(s7, 7)

# ─────────────────────────────────────────────
prs.save('presentacion_peaje.pptx')
print('PPTX generado OK')
